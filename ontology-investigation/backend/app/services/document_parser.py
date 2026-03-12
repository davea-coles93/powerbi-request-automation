"""Document Parser — AI-assisted extraction of ontology elements from documents.

Pass 1 of the two-pass ingestion architecture for unstructured content.
Unlike spreadsheets and Power BI (which use deterministic parsing), documents
require AI to understand and extract structured elements.

Two-stage approach:
  Stage 1: Text/image extraction using Python libraries (pdfplumber, python-docx, python-pptx)
  Stage 2: Haiku AI call to interpret content and extract ontology elements

Supported formats:
  - PDF (.pdf) — text extraction with page-level chunking; image pages sent via vision
  - Word (.docx) — paragraph and table extraction
  - PowerPoint (.pptx) — slide text + speaker notes extraction
  - Images (.png, .jpg, .jpeg, .gif, .bmp, .webp) — sent directly to Haiku via vision

The output matches the same staged-source format as other parsers so it feeds
directly into the cross-source AI enrichment (Pass 2, Sonnet).
"""

import base64
import io
import logging
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ── Helpers ──────────────────────────────────────────────────────────────


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx"} | IMAGE_EXTENSIONS

MAX_TEXT_CHARS = 80_000   # Truncate extracted text to stay within Haiku context
MAX_IMAGE_SIZE = 5 * 1024 * 1024  # 5 MB per image for vision API

# Chunking: if extracted text exceeds this, split into multiple Haiku calls
# and merge. Haiku 4.5 has 200K context but only 8192 output tokens, so
# sending too much in one call risks hitting the output ceiling.
CHUNK_CHAR_THRESHOLD = 15_000  # ~4K tokens — split above this
MAX_CHUNKS = 6  # Cost control: max parallel Haiku calls per document


def _slug(name: str) -> str:
    """Generate a clean snake_case ID from a name."""
    slug = name.lower().strip()
    slug = re.sub(r"[^a-z0-9_\s-]", "", slug)
    slug = re.sub(r"[\s-]+", "_", slug)
    slug = re.sub(r"_+", "_", slug)
    return slug.strip("_")


def is_image_file(filename: str) -> bool:
    """Check if filename is a supported image format."""
    return Path(filename).suffix.lower() in IMAGE_EXTENSIONS


def is_document_file(filename: str) -> bool:
    """Check if filename is a supported document format."""
    return Path(filename).suffix.lower() in DOCUMENT_EXTENSIONS


def _get_media_type(filename: str) -> str:
    """Map file extension to MIME type for vision API."""
    ext = Path(filename).suffix.lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".webp": "image/webp",
    }.get(ext, "image/png")


# ── Text Extraction ─────────────────────────────────────────────────────


def extract_text_from_pdf(data: bytes) -> list[dict]:
    """Extract text from PDF, page by page.

    Returns a list of {page: int, text: str} dicts.
    Falls back to empty text for pages that fail extraction.
    """
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed — PDF text extraction unavailable")
        return [{"page": 1, "text": "(PDF text extraction requires pdfplumber)"}]

    pages = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages):
                try:
                    text = page.extract_text() or ""
                    # Also extract tables as text
                    tables = page.extract_tables()
                    table_text = ""
                    for table in tables:
                        if table:
                            rows = []
                            for row in table:
                                cells = [str(c) if c else "" for c in row]
                                rows.append(" | ".join(cells))
                            table_text += "\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]\n"

                    full_text = text
                    if table_text:
                        full_text += "\n" + table_text

                    pages.append({"page": i + 1, "text": full_text.strip()})
                except Exception as e:
                    logger.warning(f"Failed to extract page {i + 1}: {e}")
                    pages.append({"page": i + 1, "text": ""})
    except Exception as e:
        logger.error(f"Failed to open PDF: {e}")
        return [{"page": 1, "text": f"(Failed to open PDF: {e})"}]

    return pages


def extract_text_from_docx(data: bytes) -> list[dict]:
    """Extract text from a Word document.

    Returns sections based on headings. If no headings, returns one section.
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed — DOCX extraction unavailable")
        return [{"section": "document", "text": "(DOCX extraction requires python-docx)"}]

    sections = []
    try:
        doc = Document(io.BytesIO(data))

        current_section = "Introduction"
        current_text: list[str] = []

        for para in doc.paragraphs:
            # Detect heading styles to create sections
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                # Save previous section
                if current_text:
                    sections.append({
                        "section": current_section,
                        "text": "\n".join(current_text).strip(),
                    })
                current_section = para.text.strip() or current_section
                current_text = []
            else:
                if para.text.strip():
                    current_text.append(para.text)

        # Save last section
        if current_text:
            sections.append({
                "section": current_section,
                "text": "\n".join(current_text).strip(),
            })

        # Extract tables
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                sections.append({
                    "section": f"Table {i + 1}",
                    "text": "[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]",
                })

    except Exception as e:
        logger.error(f"Failed to parse DOCX: {e}")
        return [{"section": "document", "text": f"(Failed to parse DOCX: {e})"}]

    return sections if sections else [{"section": "document", "text": "(Empty document)"}]


def extract_text_from_pptx(data: bytes) -> list[dict]:
    """Extract text and embedded images from a PowerPoint presentation.

    Returns one entry per slide with text from shapes + speaker notes.
    Slides containing embedded images (charts, screenshots, diagrams) include
    the image data so it can be sent to Haiku via vision API.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("python-pptx not installed — PPTX extraction unavailable")
        return [{"slide": 1, "text": "(PPTX extraction requires python-pptx)"}]

    slides = []
    try:
        prs = Presentation(io.BytesIO(data))

        for i, slide in enumerate(prs.slides):
            parts: list[str] = []
            slide_images: list[dict] = []

            # Extract text and images from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

                # Extract table content
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        parts.append("[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")

                # Extract embedded images (pictures, charts rendered as images)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        content_type = image.content_type or "image/png"
                        if len(img_bytes) <= MAX_IMAGE_SIZE:
                            slide_images.append({
                                "data": img_bytes,
                                "content_type": content_type,
                            })
                    except Exception:
                        pass  # Skip images we can't extract

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[SPEAKER NOTES]\n{notes}\n[/SPEAKER NOTES]")

            entry: dict = {
                "slide": i + 1,
                "text": "\n\n".join(parts).strip(),
            }
            if slide_images:
                entry["images"] = slide_images

            slides.append(entry)

    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")
        return [{"slide": 1, "text": f"(Failed to parse PPTX: {e})"}]

    return slides if slides else [{"slide": 1, "text": "(Empty presentation)"}]


def extract_images_from_pdf(data: bytes, max_images: int = 10) -> list[dict]:
    """Extract embedded images from PDF pages.

    Returns list of {page: int, data: bytes, content_type: str}.
    Limited to max_images to control token cost.
    """
    try:
        import pdfplumber
    except ImportError:
        return []

    images = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages):
                if len(images) >= max_images:
                    break
                # pdfplumber gives us image metadata; for actual extraction
                # we'd need to use the underlying PDF objects. For now, if a
                # page has images but very little text, flag it as visual.
                page_images = page.images or []
                page_text = (page.extract_text() or "").strip()

                # If the page is image-heavy and text-light, it's likely a
                # chart/diagram — we can't easily extract the raw image from
                # pdfplumber, but we note it for the AI prompt
                if len(page_images) > 0 and len(page_text) < 50:
                    images.append({
                        "page": i + 1,
                        "image_count": len(page_images),
                        "is_visual_page": True,
                    })
    except Exception:
        pass

    return images


# ── AI Extraction (Tool Use) ─────────────────────────────────────────────


EXTRACTION_SYSTEM_PROMPT = """\
You are a senior data analyst extracting structured business ontology elements from a document.

## Ontology Framework

You are building a business ontology that connects operational data creation through \
to financial analysis. The chain is:

  Attribute (raw data) → Measure (calculation) → Metric (business KPI answering a question)

Grouped by Entity (business object like Customer, Production Order), sourced from Systems \
(SAP, Excel, Power BI), viewed through Perspectives:
  - **operational**: What work is being done? Raw data, transactions, source systems.
  - **management**: How are we performing? KPIs, monitoring, aggregated views.
  - **financial**: What's the financial position? Revenue, costs, margins, profitability.

## Critical Thinking — Quality Over Quantity

**Be selective and skeptical.** Your job is NOT to extract everything that could possibly \
be an element. Your job is to extract what genuinely IS an ontology element based on \
clear evidence in the document.

Apply these filters before creating any element:

- **Entity test**: Is this a real business object that has its own attributes and participates \
in relationships? "Revenue" is not an entity — it's a measure. "Customer" IS an entity.
- **Attribute test**: Is this a specific data field that would exist as a column in a table? \
"Performance" is not an attribute — it's vague. "Performance Rating (1-5)" is an attribute.
- **Measure test**: Is there a concrete calculation described? "Costs" is not a measure. \
"Total Material Cost = SUM(Unit Cost × Quantity Issued)" IS a measure.
- **Metric test**: Does this answer a specific business question that someone actually asks? \
Don't promote every measure to a metric. Only create metrics for KPIs that drive decisions.
- **System test**: Is a specific application or data store named? Don't create systems for \
generic concepts like "database" or "reporting tool" — only for named systems (SAP, Salesforce, \
Power BI, a specific Excel workbook).

**When in doubt, leave it out.** It is much better to return 5 high-quality, well-evidenced \
elements than 20 speculative ones. The downstream enrichment pass (Sonnet) will infer \
additional elements — your job is to capture what the document explicitly states.

**Empty arrays are perfectly fine.** If a document doesn't mention any measures, return an \
empty measures array. Do not fabricate elements to fill gaps.

## What to Extract

From the document, extract:
1. **Structured elements** — entities, attributes, measures, metrics, systems, relationships, processes
2. **Business context** — rules, constraints, assumptions, thresholds, SLAs, policies mentioned
3. **Qualitative insights** — pain points, opportunities, risks, dependencies, tribal knowledge
4. **Rich descriptions** — don't just name things, capture *why* they matter and *how* they're used

## Guidelines
- Use snake_case for all IDs
- Only extract what is explicitly mentioned or clearly implied — do not invent
- If the document is a screenshot/image, describe what you see and extract visible elements
- Link attributes to entities where the relationship is clear
- Set perspective_ids based on context clues
- Capture nuance: if the document says "this report takes 3 days to produce manually", \
that's a key_insight, not just a process step
- Business rules like "revenue is only recognized when goods ship" should be captured
- If the document references existing systems (SAP, Salesforce, etc.), create system entries
- For each element, ask yourself: "What is my evidence from the document for this?"
- If you can't point to specific text or visual evidence, don't create the element
"""

# Tool schema for structured output — forces Haiku to return conformant JSON
EXTRACTION_TOOL_SCHEMA = {
    "name": "extract_ontology_elements",
    "description": "Extract structured ontology elements and business context from a document.",
    "input_schema": {
        "type": "object",
        "properties": {
            "document_type": {
                "type": "string",
                "enum": [
                    "data_dictionary", "requirements_doc", "process_doc",
                    "report_screenshot", "presentation", "architecture_doc",
                    "meeting_notes", "other",
                ],
                "description": "Classification of the document type.",
            },
            "document_summary": {
                "type": "string",
                "description": "2-3 sentence summary: what the document contains, who it's for, and what decisions it supports.",
            },
            "confidence": {
                "type": "string",
                "enum": ["high", "medium", "low"],
                "description": "Confidence in extraction quality. Low if document is unclear, heavily visual, or not business-data-oriented.",
            },
            "extraction_notes": {
                "type": "string",
                "description": "Brief notes on your extraction decisions: what you included, what you deliberately excluded and why, any elements you were uncertain about. This helps downstream reviewers understand your reasoning.",
            },
            "key_insights": {
                "type": "array",
                "description": "Qualitative observations, pain points, opportunities, risks, or tribal knowledge found in the document. These capture nuance that doesn't fit into structured elements.",
                "items": {
                    "type": "object",
                    "properties": {
                        "insight": {"type": "string", "description": "The observation or finding."},
                        "category": {
                            "type": "string",
                            "enum": ["pain_point", "opportunity", "risk", "dependency", "assumption", "tribal_knowledge", "observation"],
                            "description": "Type of insight.",
                        },
                        "relevant_entities": {
                            "type": "array", "items": {"type": "string"},
                            "description": "IDs of entities this insight relates to.",
                        },
                        "source_location": {"type": "string", "description": "Where in the document this was found (e.g., 'Page 3', 'Slide 7', 'Section: Revenue Recognition')."},
                    },
                    "required": ["insight", "category"],
                },
            },
            "business_rules": {
                "type": "array",
                "description": "Explicit or implied business rules, constraints, policies, thresholds, or SLAs mentioned in the document.",
                "items": {
                    "type": "object",
                    "properties": {
                        "rule": {"type": "string", "description": "The business rule or constraint."},
                        "context": {"type": "string", "description": "Why this rule exists or when it applies."},
                        "affects": {
                            "type": "array", "items": {"type": "string"},
                            "description": "IDs of elements (entities, measures, metrics) this rule affects.",
                        },
                    },
                    "required": ["rule"],
                },
            },
            "entities": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Display name."},
                        "description": {"type": "string", "description": "What this entity represents and why it matters in the business context. Include how it's used, not just what it is."},
                        "context_notes": {"type": "string", "description": "Additional context: ownership, lifecycle, known issues, or how this entity relates to business processes."},
                        "core_attributes": {"type": "array", "items": {"type": "string"}},
                        "lenses": {"type": "array", "items": {"type": "string"}},
                        "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                    },
                    "required": ["id", "name", "description"],
                },
            },
            "attributes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Display name."},
                        "description": {"type": "string", "description": "What this attribute captures, its business meaning, and any known data quality considerations."},
                        "context_notes": {"type": "string", "description": "Additional context: how it's collected, who enters it, known issues, manual workarounds."},
                        "entity_id": {"type": "string", "description": "Parent entity ID."},
                        "system_id": {"type": "string", "description": "Source system ID if known."},
                        "source_table": {"type": "string", "description": "Source table/dataset if mentioned."},
                        "source_column": {"type": "string", "description": "Source column/field if mentioned."},
                        "data_type": {"type": "string", "enum": ["string", "number", "datetime", "boolean", "currency", "percentage"]},
                        "perspective_ids": {"type": "array", "items": {"type": "string"}},
                        "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                    },
                    "required": ["id", "name", "description"],
                },
            },
            "measures": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Display name."},
                        "description": {"type": "string", "description": "What this measure calculates, its business purpose, and how it's typically used in analysis."},
                        "context_notes": {"type": "string", "description": "Additional context: caveats, known limitations, manual adjustments applied, seasonality effects."},
                        "logic": {"type": "string", "description": "Plain English calculation logic."},
                        "formula": {"type": "string", "description": "Formal formula if mentioned (DAX, SQL, Excel formula)."},
                        "input_attribute_ids": {"type": "array", "items": {"type": "string"}},
                        "input_measure_ids": {"type": "array", "items": {"type": "string"}},
                        "perspective_ids": {"type": "array", "items": {"type": "string"}},
                        "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                    },
                    "required": ["id", "name", "description"],
                },
            },
            "metrics": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Display name."},
                        "description": {"type": "string", "description": "What this KPI measures, why leadership cares about it, and what actions it drives."},
                        "context_notes": {"type": "string", "description": "Additional context: targets/thresholds, reporting frequency, who reviews it, trend expectations."},
                        "business_question": {"type": "string", "description": "The business question this metric answers."},
                        "target": {"type": "string", "description": "Target value or threshold if mentioned (e.g., '>95%', '<3 days')."},
                        "frequency": {"type": "string", "description": "How often this metric is reviewed if mentioned (daily, weekly, monthly, quarterly)."},
                        "calculated_by_measure_ids": {"type": "array", "items": {"type": "string"}},
                        "perspective_ids": {"type": "array", "items": {"type": "string"}},
                        "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                    },
                    "required": ["id", "name", "description"],
                },
            },
            "systems": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Display name."},
                        "type": {"type": "string", "enum": ["ERP", "BI", "Spreadsheet", "Database", "CRM", "WMS", "MES", "HCM", "Other"]},
                        "vendor": {"type": "string", "description": "Vendor name if known."},
                        "integration_status": {"type": "string", "enum": ["Connected", "Manual Extract", "API", "Unknown"]},
                        "context_notes": {"type": "string", "description": "Additional context: known limitations, who manages it, shadow system concerns, integration pain points."},
                        "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                    },
                    "required": ["id", "name", "type"],
                },
            },
            "relationships": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string"},
                        "name": {"type": "string", "description": "Descriptive relationship name."},
                        "from_entity_id": {"type": "string"},
                        "to_entity_id": {"type": "string"},
                        "relationship_type": {"type": "string", "enum": ["many-to-one", "one-to-many", "many-to-many", "one-to-one"]},
                        "description": {"type": "string", "description": "How and why these entities relate."},
                    },
                    "required": ["id", "from_entity_id", "to_entity_id", "relationship_type"],
                },
            },
            "processes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "string", "description": "snake_case identifier."},
                        "name": {"type": "string", "description": "Process name."},
                        "description": {"type": "string", "description": "What this process accomplishes, its business purpose, and who's involved."},
                        "frequency": {"type": "string", "description": "How often this process runs if mentioned."},
                        "pain_points": {"type": "array", "items": {"type": "string"}, "description": "Known pain points or inefficiencies in this process."},
                        "steps": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "id": {"type": "string"},
                                    "sequence": {"type": "integer"},
                                    "name": {"type": "string"},
                                    "description": {"type": "string", "description": "What happens in this step, including any manual effort or system interactions."},
                                    "perspective_id": {"type": "string"},
                                    "systems_used_ids": {"type": "array", "items": {"type": "string"}},
                                    "manual_effort_percentage": {"type": "integer", "description": "Estimated manual effort 0-100 if mentioned."},
                                    "depends_on_step_ids": {"type": "array", "items": {"type": "string"}},
                                    "state": {"type": "string", "enum": ["as-is", "to-be", "gap"]},
                                },
                                "required": ["id", "sequence", "name"],
                            },
                        },
                    },
                    "required": ["id", "name", "description"],
                },
            },
        },
        "required": ["document_type", "document_summary", "confidence"],
    },
}


MAX_VISION_IMAGES = 8  # Limit embedded images sent to Haiku (cost control)


def _build_content_blocks(
    extracted_content: list[dict],
    doc_type: str,
    filename: str,
    image_data: Optional[bytes] = None,
    image_filename: Optional[str] = None,
) -> list[dict]:
    """Build the message content blocks for the Haiku call.

    For text documents: text content with page/section markers.
    For images: vision content block with base64-encoded image.
    For PPTX/PDF with embedded images: mixed text + vision blocks.
    """
    blocks: list[dict] = []

    # If this is a standalone image file, send it via vision
    if image_data and image_filename:
        media_type = _get_media_type(image_filename)
        b64 = base64.b64encode(image_data).decode("utf-8")
        blocks.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": b64,
            },
        })
        blocks.append({
            "type": "text",
            "text": f"This is an image file: {filename}. Extract any business ontology elements visible in this image.",
        })
        return blocks

    # For text documents, format the extracted content
    # Collect embedded images to append as vision blocks
    text_parts = [f"Document: {filename} (type: {doc_type})\n"]
    embedded_images: list[dict] = []

    total_chars = 0
    for chunk in extracted_content:
        # Format based on content type
        if "page" in chunk:
            header = f"\n--- Page {chunk['page']} ---\n"
        elif "slide" in chunk:
            header = f"\n--- Slide {chunk['slide']} ---\n"
        elif "section" in chunk:
            header = f"\n--- {chunk['section']} ---\n"
        else:
            header = "\n---\n"

        text = chunk.get("text", "")

        # Collect embedded images from slides (PPTX)
        for img in chunk.get("images", []):
            if len(embedded_images) < MAX_VISION_IMAGES:
                label = f"Slide {chunk.get('slide', '?')}"
                embedded_images.append({
                    "data": img["data"],
                    "content_type": img.get("content_type", "image/png"),
                    "label": label,
                })

        if not text.strip():
            # Even if no text, note that images were found
            if chunk.get("images"):
                slide_num = chunk.get("slide", "?")
                text_parts.append(
                    f"{header}[This slide contains {len(chunk['images'])} embedded image(s) — see below]"
                )
            continue

        # Truncate if we're getting too long
        remaining = MAX_TEXT_CHARS - total_chars
        if remaining <= 0:
            text_parts.append("\n(... content truncated for length ...)")
            break

        if len(text) > remaining:
            text = text[:remaining] + "\n(... truncated ...)"

        img_note = ""
        if chunk.get("images"):
            img_note = f"\n[+ {len(chunk['images'])} embedded image(s) — see below]"

        text_parts.append(header + text + img_note)
        total_chars += len(header) + len(text) + len(img_note)

    blocks.append({"type": "text", "text": "\n".join(text_parts)})

    # Append embedded images as vision blocks (PPTX slides, etc.)
    if embedded_images:
        blocks.append({
            "type": "text",
            "text": (
                f"\n\nThe following {len(embedded_images)} embedded image(s) were extracted "
                f"from the document. These may contain charts, diagrams, screenshots, or "
                f"other visual content with business data. Examine each and extract any "
                f"ontology elements visible."
            ),
        })
        for img_info in embedded_images:
            b64 = base64.b64encode(img_info["data"]).decode("utf-8")
            blocks.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": img_info["content_type"],
                    "data": b64,
                },
            })
            blocks.append({
                "type": "text",
                "text": f"[Embedded image from {img_info['label']}]",
            })

    return blocks


def _chunk_content(extracted_content: list[dict]) -> list[list[dict]]:
    """Split extracted content into chunks if it exceeds the threshold.

    Each chunk is a list of page/section/slide dicts. We split at natural
    boundaries (pages, sections, slides) rather than mid-text.
    """
    # Calculate total text size
    total_chars = sum(len(c.get("text", "")) for c in extracted_content)

    if total_chars <= CHUNK_CHAR_THRESHOLD:
        return [extracted_content]

    # Target chunk size — split evenly but respect MAX_CHUNKS
    n_chunks = min(MAX_CHUNKS, max(2, total_chars // CHUNK_CHAR_THRESHOLD + 1))
    target_per_chunk = total_chars // n_chunks

    chunks: list[list[dict]] = []
    current_chunk: list[dict] = []
    current_size = 0

    for item in extracted_content:
        item_size = len(item.get("text", ""))
        current_chunk.append(item)
        current_size += item_size

        if current_size >= target_per_chunk and len(chunks) < n_chunks - 1:
            chunks.append(current_chunk)
            current_chunk = []
            current_size = 0

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def _merge_chunk_results(chunk_results: list[dict]) -> dict:
    """Merge extraction results from multiple chunks into one result.

    Combines all array fields, takes the best document_type/confidence,
    and concatenates summaries.
    """
    if len(chunk_results) == 1:
        return chunk_results[0]

    merged = {
        "document_type": chunk_results[0].get("document_type", "other"),
        "document_summary": "",
        "confidence": "high",
        "key_insights": [],
        "business_rules": [],
        "entities": [],
        "attributes": [],
        "measures": [],
        "metrics": [],
        "systems": [],
        "relationships": [],
        "processes": [],
    }

    summaries = []
    extraction_notes_parts = []
    confidence_rank = {"high": 3, "medium": 2, "low": 1}
    worst_confidence = 3

    # Collect all IDs to detect duplicates across chunks
    seen_ids: dict[str, set] = {
        "entities": set(), "attributes": set(), "measures": set(),
        "metrics": set(), "systems": set(), "relationships": set(),
        "processes": set(),
    }

    total_input_tokens = 0
    total_output_tokens = 0

    for result in chunk_results:
        # Summaries
        s = result.get("document_summary", "")
        if s:
            summaries.append(s)

        # Extraction notes
        en = result.get("extraction_notes", "")
        if en:
            extraction_notes_parts.append(en)

        # Confidence — take the worst
        c = result.get("confidence", "medium")
        worst_confidence = min(worst_confidence, confidence_rank.get(c, 2))

        # Insights and rules — just concatenate (Pass 2 dedupes)
        merged["key_insights"].extend(result.get("key_insights", []))
        merged["business_rules"].extend(result.get("business_rules", []))

        # Structural elements — skip exact ID duplicates
        for field in seen_ids:
            for item in result.get(field, []):
                item_id = item.get("id", "")
                if item_id and item_id in seen_ids[field]:
                    continue
                if item_id:
                    seen_ids[field].add(item_id)
                merged[field].append(item)

        # Token usage
        tu = result.get("token_usage", {})
        total_input_tokens += tu.get("input_tokens", 0)
        total_output_tokens += tu.get("output_tokens", 0)

    # Best summary — use the first chunk's (it has the overview context)
    merged["document_summary"] = summaries[0] if summaries else ""
    merged["extraction_notes"] = " | ".join(extraction_notes_parts) if extraction_notes_parts else ""

    merged["confidence"] = {3: "high", 2: "medium", 1: "low"}.get(worst_confidence, "medium")

    merged["token_usage"] = {
        "model": chunk_results[0].get("token_usage", {}).get("model", ""),
        "input_tokens": total_input_tokens,
        "output_tokens": total_output_tokens,
        "chunks": len(chunk_results),
    }

    return merged


async def _extract_single_chunk(
    client,
    model: str,
    system: str,
    content: list[dict],
    filename: str,
    chunk_label: str = "",
) -> dict:
    """Make a single Haiku tool-use call and return the parsed result."""
    messages = [{"role": "user", "content": content}]

    response = await client.messages.create(
        model=model,
        max_tokens=8192,
        system=system,
        messages=messages,
        tools=[EXTRACTION_TOOL_SCHEMA],
        tool_choice={"type": "tool", "name": "extract_ontology_elements"},
    )

    parsed = None
    for block in response.content:
        if block.type == "tool_use" and block.name == "extract_ontology_elements":
            parsed = block.input
            break

    if not parsed:
        logger.warning(f"No tool_use block in response for {filename} {chunk_label}")
        return _empty_result(filename, f"(AI did not return structured data {chunk_label})")

    # Extract token usage
    token_usage = {}
    if hasattr(response, "usage") and response.usage:
        token_usage = {
            "model": model,
            "input_tokens": response.usage.input_tokens,
            "output_tokens": response.usage.output_tokens,
        }
        logger.info(
            f"Extraction tokens for {filename} {chunk_label}: "
            f"input={response.usage.input_tokens}, output={response.usage.output_tokens}"
        )

    return {
        "source_type": "document",
        "source_name": Path(filename).stem,
        "document_type": parsed.get("document_type", "other"),
        "document_summary": parsed.get("document_summary", ""),
        "confidence": parsed.get("confidence", "medium"),
        "extraction_notes": parsed.get("extraction_notes", ""),
        "key_insights": parsed.get("key_insights", []),
        "business_rules": parsed.get("business_rules", []),
        "entities": parsed.get("entities", []),
        "attributes": parsed.get("attributes", []),
        "measures": parsed.get("measures", []),
        "metrics": parsed.get("metrics", []),
        "systems": parsed.get("systems", []),
        "relationships": parsed.get("relationships", []),
        "processes": parsed.get("processes", []),
        "perspectives": [],
        "warnings": [],
        "token_usage": token_usage,
    }


async def extract_elements_with_ai(
    extracted_content: list[dict],
    doc_type: str,
    filename: str,
    ontology_context: str = "",
    image_data: Optional[bytes] = None,
    image_filename: Optional[str] = None,
) -> dict:
    """Call Haiku to extract structured ontology elements from document content.

    Uses Anthropic tool use to guarantee structured JSON output. For large
    documents, automatically chunks content and runs parallel extraction
    calls, then merges results. Haiku 4.5 has 8192 max output tokens, so
    chunking ensures we get full extraction from dense documents.

    Args:
        extracted_content: List of page/section/slide dicts with text
        doc_type: Document format (pdf, docx, pptx, image)
        filename: Original filename
        ontology_context: Current ontology state summary (from OntologyContextBuilder)
        image_data: Raw image bytes (for image files)
        image_filename: Image filename (for MIME type detection)

    Returns:
        Dict matching the staged source format (entities, attributes, measures, etc.)
    """
    import asyncio
    import os

    try:
        import anthropic
    except ImportError:
        raise RuntimeError("anthropic package required for document extraction")

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise RuntimeError("ANTHROPIC_API_KEY required for document extraction")

    from ..config import AI_EXTRACTION_MODEL

    client = anthropic.AsyncAnthropic(api_key=api_key)

    # Build system prompt with ontology context
    system = EXTRACTION_SYSTEM_PROMPT
    if ontology_context and "EMPTY" not in ontology_context:
        system += (
            "\n\n## Existing Ontology (reference these IDs where relevant — "
            "don't re-create elements that already exist)\n\n"
            + ontology_context
        )

    # For standalone images, no chunking — single call with vision
    if image_data and image_filename:
        content = _build_content_blocks(
            [], "image", filename,
            image_data=image_data, image_filename=image_filename,
        )
        try:
            return await _extract_single_chunk(
                client, AI_EXTRACTION_MODEL, system, content, filename,
            )
        except Exception as e:
            logger.error(f"Haiku extraction failed for {filename}: {e}")
            raise RuntimeError(f"AI extraction failed: {e}")

    # Chunk large documents for parallel extraction
    chunks = _chunk_content(extracted_content)

    if len(chunks) == 1:
        # Small document — single call
        content = _build_content_blocks(extracted_content, doc_type, filename)
        try:
            return await _extract_single_chunk(
                client, AI_EXTRACTION_MODEL, system, content, filename,
            )
        except Exception as e:
            logger.error(f"Haiku extraction failed for {filename}: {e}")
            raise RuntimeError(f"AI extraction failed: {e}")

    # Large document — parallel chunk extraction
    logger.info(f"Chunking {filename} into {len(chunks)} parts for parallel extraction")

    async def extract_chunk(i: int, chunk: list[dict]) -> dict:
        label = f"(chunk {i + 1}/{len(chunks)})"
        content = _build_content_blocks(chunk, doc_type, f"{filename} {label}")
        return await _extract_single_chunk(
            client, AI_EXTRACTION_MODEL, system, content, filename, label,
        )

    try:
        chunk_results = await asyncio.gather(
            *(extract_chunk(i, chunk) for i, chunk in enumerate(chunks)),
            return_exceptions=True,
        )

        # Filter out failures
        successful = []
        for i, result in enumerate(chunk_results):
            if isinstance(result, Exception):
                logger.warning(f"Chunk {i + 1}/{len(chunks)} failed for {filename}: {result}")
            else:
                successful.append(result)

        if not successful:
            raise RuntimeError("All chunks failed extraction")

        merged = _merge_chunk_results(successful)
        # Re-attach source metadata
        merged["source_type"] = "document"
        merged["source_name"] = Path(filename).stem
        merged["perspectives"] = []
        if len(successful) < len(chunks):
            merged.setdefault("warnings", []).append(
                f"{len(chunks) - len(successful)}/{len(chunks)} chunks failed"
            )
        else:
            merged.setdefault("warnings", [])

        return merged

    except Exception as e:
        logger.error(f"Haiku extraction failed for {filename}: {e}")
        raise RuntimeError(f"AI extraction failed: {e}")


def _empty_result(filename: str, warning: str) -> dict:
    """Return an empty result with a warning."""
    return {
        "source_type": "document",
        "source_name": Path(filename).stem,
        "document_type": "other",
        "document_summary": "",
        "confidence": "low",
        "entities": [],
        "attributes": [],
        "measures": [],
        "metrics": [],
        "systems": [],
        "relationships": [],
        "processes": [],
        "perspectives": [],
        "warnings": [warning],
    }


# ── Main Parse Entry Point ──────────────────────────────────────────────


async def parse_document(
    data: bytes,
    filename: str,
    ontology_context: str = "",
) -> dict:
    """Parse a document file and extract ontology elements using Haiku.

    This is the main entry point for document ingestion. It:
    1. Detects the file type
    2. Extracts text/content using the appropriate library
    3. Sends to Haiku for structured element extraction
    4. Returns the result in the standard staged source format

    Args:
        data: Raw file bytes
        filename: Original filename (used for type detection)
        ontology_context: Current ontology summary for AI context

    Returns:
        Dict matching the staged source format
    """
    ext = Path(filename).suffix.lower()

    # Route to appropriate extractor
    if ext == ".pdf":
        extracted = extract_text_from_pdf(data)
        doc_type = "pdf"
        result = await extract_elements_with_ai(
            extracted, doc_type, filename,
            ontology_context=ontology_context,
        )

    elif ext == ".docx":
        extracted = extract_text_from_docx(data)
        doc_type = "docx"
        result = await extract_elements_with_ai(
            extracted, doc_type, filename,
            ontology_context=ontology_context,
        )

    elif ext == ".pptx":
        extracted = extract_text_from_pptx(data)
        doc_type = "pptx"
        result = await extract_elements_with_ai(
            extracted, doc_type, filename,
            ontology_context=ontology_context,
        )

    elif ext in IMAGE_EXTENSIONS:
        # Images go directly to Haiku via vision
        if len(data) > MAX_IMAGE_SIZE:
            return _empty_result(filename, f"Image too large (max {MAX_IMAGE_SIZE // (1024*1024)}MB)")

        result = await extract_elements_with_ai(
            [], "image", filename,
            ontology_context=ontology_context,
            image_data=data,
            image_filename=filename,
        )

    else:
        return _empty_result(filename, f"Unsupported file type: {ext}")

    return result
